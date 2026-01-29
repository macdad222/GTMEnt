"""PowerPoint deck generator for playbooks."""

from typing import Optional
from datetime import datetime
import io

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from src.playbook.models import Playbook, PlaybookSection


class PPTXGenerator:
    """
    Generate PowerPoint decks from playbooks.

    Produces consulting-grade slides with:
    - Title slide
    - Table of contents
    - Section slides with narrative + key points
    - Appendix with citations and assumptions
    """

    # Brand colors
    BRAND_BLUE = RgbColor(0, 102, 204)  # #0066cc
    ACCENT_ORANGE = RgbColor(236, 118, 18)  # #ec7612
    DARK_BG = RgbColor(23, 32, 51)  # #172033
    TEXT_WHITE = RgbColor(255, 255, 255)
    TEXT_GRAY = RgbColor(148, 163, 184)

    def __init__(self):
        self.prs: Optional[Presentation] = None

    def generate(self, playbook: Playbook, include_appendix: bool = True) -> bytes:
        """
        Generate a PPTX file from a playbook.

        Args:
            playbook: The playbook to convert
            include_appendix: Whether to include citations/assumptions appendix

        Returns:
            PPTX file as bytes
        """
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)

        # Generate slides
        self._add_title_slide(playbook)
        self._add_toc_slide(playbook)

        for section in playbook.sections:
            if section.section_type != "appendix":
                self._add_section_slide(section)

        if include_appendix:
            self._add_appendix_slide(playbook)

        # Export to bytes
        output = io.BytesIO()
        self.prs.save(output)
        output.seek(0)
        return output.read()

    def _add_title_slide(self, playbook: Playbook) -> None:
        """Add the title slide."""
        slide_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(slide_layout)

        # Background
        self._set_slide_background(slide, self.DARK_BG)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(1.5))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = playbook.name
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_WHITE
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.83), Inches(0.8))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = playbook.description
        p.font.size = Pt(20)
        p.font.color.rgb = self.TEXT_GRAY
        p.alignment = PP_ALIGN.LEFT

        # Date
        date_box = slide.shapes.add_textbox(Inches(0.75), Inches(6.5), Inches(5), Inches(0.4))
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Generated: {datetime.utcnow().strftime('%B %d, %Y')}"
        p.font.size = Pt(14)
        p.font.color.rgb = self.TEXT_GRAY

        # Branding
        brand_box = slide.shapes.add_textbox(Inches(10), Inches(6.5), Inches(3), Inches(0.4))
        tf = brand_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Comcast Business"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = self.BRAND_BLUE
        p.alignment = PP_ALIGN.RIGHT

    def _add_toc_slide(self, playbook: Playbook) -> None:
        """Add table of contents slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, self.DARK_BG)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Table of Contents"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_WHITE

        # TOC items
        y_offset = 1.6
        for i, section in enumerate(playbook.sections, 1):
            item_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_offset), Inches(11.83), Inches(0.5))
            tf = item_box.text_frame
            p = tf.paragraphs[0]
            p.text = f"{i}. {section.title}"
            p.font.size = Pt(18)
            p.font.color.rgb = self.TEXT_WHITE
            y_offset += 0.6

    def _add_section_slide(self, section: PlaybookSection) -> None:
        """Add a section content slide."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, self.DARK_BG)

        # Section title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_WHITE

        # Narrative
        if section.narrative:
            narrative_box = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(1.5))
            tf = narrative_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = section.narrative
            p.font.size = Pt(16)
            p.font.color.rgb = self.TEXT_GRAY
            p.line_spacing = 1.3

        # Key points
        if section.key_points:
            y_offset = 3.5
            for point in section.key_points:
                point_box = slide.shapes.add_textbox(Inches(1), Inches(y_offset), Inches(11.33), Inches(0.4))
                tf = point_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {point}"
                p.font.size = Pt(14)
                p.font.color.rgb = self.TEXT_WHITE
                y_offset += 0.5

    def _add_appendix_slide(self, playbook: Playbook) -> None:
        """Add appendix with citations and assumptions."""
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        self._set_slide_background(slide, self.DARK_BG)

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.5), Inches(11.83), Inches(0.8))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Appendix: Sources & Assumptions"
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self.TEXT_WHITE

        # Citations
        citations = playbook.get_all_citations()
        assumptions = playbook.get_all_assumptions()

        y_offset = 1.5

        if citations:
            cit_title = slide.shapes.add_textbox(Inches(0.75), Inches(y_offset), Inches(5), Inches(0.4))
            tf = cit_title.text_frame
            p = tf.paragraphs[0]
            p.text = "Sources"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.ACCENT_ORANGE
            y_offset += 0.5

            for cit in citations[:5]:  # Limit to 5
                cit_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_offset), Inches(11.83), Inches(0.4))
                tf = cit_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {cit.to_footnote()}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.TEXT_GRAY
                y_offset += 0.4

        if assumptions:
            y_offset += 0.3
            asm_title = slide.shapes.add_textbox(Inches(0.75), Inches(y_offset), Inches(5), Inches(0.4))
            tf = asm_title.text_frame
            p = tf.paragraphs[0]
            p.text = "Key Assumptions"
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = self.ACCENT_ORANGE
            y_offset += 0.5

            for asm in assumptions[:5]:  # Limit to 5
                asm_box = slide.shapes.add_textbox(Inches(0.75), Inches(y_offset), Inches(11.83), Inches(0.4))
                tf = asm_box.text_frame
                p = tf.paragraphs[0]
                p.text = f"• {asm.description}: {asm.value}"
                p.font.size = Pt(11)
                p.font.color.rgb = self.TEXT_GRAY
                y_offset += 0.4

    def _set_slide_background(self, slide, color: RgbColor) -> None:
        """Set slide background color."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

